from app import app, db, Presentation
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import os
from pptx import Presentation as PPTPresentation
from PyPDF2 import PdfReader

class NewPresentationHandler(FileSystemEventHandler):
    def on_created(self, event):
        if event.is_directory:
            return
        if event.src_path.endswith(('.pptx', '.pdf')):
            print(event.src_path)
            self.process_file(event.src_path)
            print("New presentation detected and processed: ", event.src_path)

    def process_file(self, path):
        with app.app_context():
            if path.endswith('.pptx'):
                self.process_presentation(path)
            elif path.endswith('.pdf'):
                self.process_pdf(path)

    def process_presentation(self, path): 
        ppt = PPTPresentation(path)
        title = os.path.basename(path)
        existing_entries = Presentation.query.filter_by(title=title).all()
        existing_page_numbers = set(entry.page_number for entry in existing_entries)
        for i, slide in enumerate(ppt.slides):
            if i+1 not in existing_page_numbers:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        keyword = shape.text
                        new_entry = Presentation(title=title, page_number=i+1, keyword=keyword, local_link='file/' + title)
                        db.session.add(new_entry)
        db.session.commit()
        print("Data added to database from presentation: ", title)

    def process_pdf(self, path):
        with open(path, 'rb') as file:
            pdf = PdfReader(file)
            title = os.path.basename(path)
            existing_entries = Presentation.query.filter_by(title=title).all()
            existing_page_numbers = set(entry.page_number for entry in existing_entries)
            for i, page in enumerate(pdf.pages):
                if i+1 not in existing_page_numbers:
                    text = page.extract_text()
                    for keyword in text.split('\n'):
                        if keyword.strip():
                            new_entry = Presentation(title=title, page_number=i+1, keyword=keyword, local_link='file/' + title)
                            db.session.add(new_entry)
        db.session.commit()
        print("Data added to database from PDF: ", title)

if __name__ == "__main__":
    path = 'D:/Magang/latihan/static/file'
    event_handler = NewPresentationHandler()
    observer = Observer()
    observer.schedule(event_handler, path, recursive=False)
    observer.start()
    try:
        while True:
            pass
    except KeyboardInterrupt:
        observer.stop()
    observer.join()
